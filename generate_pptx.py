"""
ADTU Academic Presentation Generator
Jira for Students – A Collaborative Academic Task Management System
BSc-IT (MAIS) 6th Semester | Assam down town University
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import pptx.util as util
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── COLOR PALETTE ──────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)
DARK_BLUE  = RGBColor(0x00, 0x32, 0x6E)   # AdtU navy
MID_BLUE   = RGBColor(0x00, 0x5B, 0xAA)   # accent blue
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)   # fill tint
GOLD       = RGBColor(0xA8, 0x7C, 0x1C)   # university gold
GRAY_TXT   = RGBColor(0x33, 0x33, 0x33)
GRAY_LIGHT = RGBColor(0xF2, 0xF5, 0xF9)
DIVIDER    = RGBColor(0x00, 0x32, 0x6E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]   # completely blank

# ── HELPERS ────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, l, t, w, h,
                font_size=Pt(14), bold=False, italic=False,
                color=BLACK, align=PP_ALIGN.LEFT,
                word_wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    txBox.word_wrap = word_wrap
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox

def add_slide_number(slide, num):
    """Bottom-right slide number"""
    add_textbox(slide, str(num),
                Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.3),
                font_size=Pt(9), color=GRAY_TXT, align=PP_ALIGN.RIGHT)

def add_header_bar(slide, chapter_tag=""):
    """Top navy bar with university name + optional chapter tag"""
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.55), fill_rgb=DARK_BLUE)
    add_textbox(slide, "Assam down town University, Guwahati",
                Inches(0.2), Inches(0.07), Inches(8), Inches(0.4),
                font_size=Pt(10), bold=True, color=WHITE, font_name="Calibri")
    if chapter_tag:
        add_textbox(slide, chapter_tag,
                    Inches(9), Inches(0.07), Inches(4.1), Inches(0.4),
                    font_size=Pt(9), color=LIGHT_BLUE,
                    align=PP_ALIGN.RIGHT, font_name="Calibri")

def add_footer_bar(slide):
    """Bottom gold accent line"""
    add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), fill_rgb=GOLD)

def add_slide_title(slide, title_text, subtitle_text="", top=Inches(0.75)):
    """Standard content slide title block"""
    add_textbox(slide, title_text,
                Inches(0.5), top, Inches(12), Inches(0.55),
                font_size=Pt(22), bold=True, color=DARK_BLUE, font_name="Calibri")
    # underline rule
    add_rect(slide, Inches(0.5), top + Inches(0.58), Inches(12), Pt(2), fill_rgb=GOLD)
    if subtitle_text:
        add_textbox(slide, subtitle_text,
                    Inches(0.5), top + Inches(0.65), Inches(12), Inches(0.35),
                    font_size=Pt(12), italic=True, color=MID_BLUE, font_name="Calibri")

def add_bullet_block(slide, items, left=Inches(0.6), top=Inches(1.6),
                     width=Inches(12.1), font_size=Pt(14),
                     color=GRAY_TXT, bullet_color=MID_BLUE, line_spacing=Pt(6)):
    """Add a list of bullet strings; supports nested (tuple: (main, [subs]))"""
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(left, top, width, Inches(5.4))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = line_spacing
        if isinstance(item, tuple):
            main_text, subs = item
        else:
            main_text, subs = item, []

        run = p.add_run()
        run.text = f"•  {main_text}"
        run.font.size = font_size
        run.font.color.rgb = color
        run.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT

        for sub in subs:
            sp = tf.add_paragraph()
            sp.space_before = Pt(2)
            sr = sp.add_run()
            sr.text = f"      ‒  {sub}"
            sr.font.size = Pt(font_size.pt - 1.5)
            sr.font.color.rgb = GRAY_TXT
            sr.font.name = "Calibri"
            sp.alignment = PP_ALIGN.LEFT


def add_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)

# Full navy top band
add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), fill_rgb=DARK_BLUE)
add_textbox(slide, "ASSAM DOWN TOWN UNIVERSITY, GUWAHATI",
            Inches(0.3), Inches(0.22), Inches(12.7), Inches(0.55),
            font_size=Pt(15), bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font_name="Calibri")

# Gold band separator
add_rect(slide, 0, Inches(1.2), SLIDE_W, Inches(0.08), fill_rgb=GOLD)

# Project title
add_textbox(slide, "Jira for Students",
            Inches(0.5), Inches(1.55), Inches(12.3), Inches(0.75),
            font_size=Pt(34), bold=True, color=DARK_BLUE,
            align=PP_ALIGN.CENTER, font_name="Calibri")
add_textbox(slide, "A Collaborative Academic Task Management System",
            Inches(0.5), Inches(2.32), Inches(12.3), Inches(0.5),
            font_size=Pt(18), italic=True, color=MID_BLUE,
            align=PP_ALIGN.CENTER, font_name="Calibri")

# Divider
add_rect(slide, Inches(3.5), Inches(2.95), Inches(6.3), Pt(1.5), fill_rgb=GOLD)

# Submitted by block
info_top = Inches(3.1)
info_lines = [
    ("Submitted By:", Pt(12), True,  DARK_BLUE),
    ("[Student Name 1]  |  Roll No: [XXXX]", Pt(13), False, GRAY_TXT),
    ("[Student Name 2]  |  Roll No: [XXXX]", Pt(13), False, GRAY_TXT),
    ("", Pt(10), False, GRAY_TXT),
    ("Course: BSc-IT (MAIS) — 6th Semester", Pt(13), False, GRAY_TXT),
    ("Project Guide: [Guide Name]", Pt(13), False, GRAY_TXT),
    ("Academic Year: 2025–2026", Pt(12), False, GRAY_TXT),
]
y = info_top
for text, sz, bd, clr in info_lines:
    add_textbox(slide, text, Inches(0.5), y, Inches(12.3), Inches(0.35),
                font_size=sz, bold=bd, color=clr,
                align=PP_ALIGN.CENTER, font_name="Calibri")
    y += Inches(0.33)

# Bottom gold bar
add_rect(slide, 0, Inches(7.25), SLIDE_W, Inches(0.25), fill_rgb=GOLD)
add_slide_number(slide, 1)

add_notes(slide, ("This is the title slide of the major project presentation.\n"
    "Introduce the project title 'Jira for Students', the course details, student names, "
    "roll numbers, and guide name. Briefly greet the audience and state that you will be "
    "presenting the final-year project developed for BSc-IT (MAIS) 6th Semester at "
    "Assam down town University."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ABSTRACT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 2)
add_slide_title(slide, "Abstract")

body = (
    "This project, 'Jira for Students', is a cloud-based academic collaboration platform "
    "developed for Assam down town University (AdtU). It replaces informal communication "
    "channels such as WhatsApp groups with a structured, role-based academic workspace.\n\n"
    "Problem: Students and faculty lack a unified digital platform for subject-specific "
    "communication, assignment tracking, file sharing, and announcements — leading to "
    "scattered information and missed deadlines.\n\n"
    "Solution: A full-stack web application built with React, Node.js, Socket.IO, PostgreSQL "
    "(Supabase), and Cloudinary that organises academic activities into dedicated subject "
    "channels with real-time messaging, assignment management, notes repository, and "
    "role-based access control.\n\n"
    "Outcome: A production-ready platform deployed on Vercel (frontend) and Render (backend) "
    "that enables seamless, structured academic collaboration between students, faculty, and "
    "class representatives."
)
add_textbox(slide, body, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.5),
            font_size=Pt(13.5), color=GRAY_TXT, word_wrap=True, font_name="Calibri")

add_notes(slide, ("Explain the abstract in 2–3 minutes.\n"
    "Three parts: (1) The problem — students rely on WhatsApp which is unstructured. "
    "(2) The solution — a Jira-inspired platform with dedicated channels per subject. "
    "(3) The outcome — a fully deployed, role-based collaborative system."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TABLE OF CONTENTS / INDEX
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 3)
add_slide_title(slide, "Index / Table of Contents")

chapters = [
    ("Chapter 1", "Introduction, Problem Statement, Objectives",   "Slides 4–6"),
    ("Chapter 2", "System Overview — Existing & Proposed Systems",  "Slides 7–9"),
    ("Chapter 3", "System Design — Architecture, DB, Workflow",     "Slides 10–12"),
    ("Chapter 4", "Implementation — Technologies, UI, Methodology", "Slides 13–15"),
    ("Chapter 5", "Results & Discussion",                           "Slides 16–17"),
    ("Chapter 6", "Conclusion & Future Scope",                      "Slides 18–19"),
    ("",          "References & Thank You",                         "Slide 20"),
]

row_h = Inches(0.6)
top   = Inches(1.55)
for i, (ch, title, pg) in enumerate(chapters):
    row_top = top + i * row_h
    bg = GRAY_LIGHT if i % 2 == 0 else WHITE
    add_rect(slide, Inches(0.4), row_top, Inches(12.5), row_h - Pt(2), fill_rgb=bg,
             line_rgb=LIGHT_BLUE, line_width=Pt(0.5))
    if ch:
        add_textbox(slide, ch, Inches(0.5), row_top + Inches(0.12),
                    Inches(1.6), Inches(0.4),
                    font_size=Pt(12), bold=True, color=DARK_BLUE, font_name="Calibri")
    add_textbox(slide, title, Inches(2.2), row_top + Inches(0.12),
                Inches(8.5), Inches(0.4),
                font_size=Pt(12), color=GRAY_TXT, font_name="Calibri")
    add_textbox(slide, pg, Inches(10.9), row_top + Inches(0.12),
                Inches(1.8), Inches(0.4),
                font_size=Pt(11), italic=True, color=MID_BLUE,
                align=PP_ALIGN.RIGHT, font_name="Calibri")

add_notes(slide, ("Walk the audience through the overall structure of the presentation.\n"
    "This slide gives a quick map of all six chapters. Mention that the project covers "
    "introduction, design, implementation, and evaluation — following a standard academic "
    "software engineering report format."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — INTRODUCTION (Chapter 1)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 1: Introduction")
add_footer_bar(slide)
add_slide_number(slide, 4)
add_slide_title(slide, "Introduction")

bullets = [
    ("Background",
     ["Academic collaboration in universities increasingly relies on digital tools.",
      "At AdtU, student-faculty communication happens largely through informal channels.",
      "No centralised platform exists for subject-specific task and resource management."]),
    ("Need for the System",
     ["Students miss assignments due to buried messages in group chats.",
      "Files and notes lack proper version control or searchability.",
      "Faculty cannot track submissions or send targeted announcements efficiently.",
      "A Jira-inspired workspace tailored for academic use is the logical solution."]),
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(13.5))

add_notes(slide, ("Introduce the context of the problem.\n"
    "Highlight that while tools like Jira exist for software teams, no equivalent exists for "
    "academic student groups. Explain how this motivated the development of a purpose-built "
    "platform for university students at AdtU."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 1: Introduction")
add_footer_bar(slide)
add_slide_number(slide, 5)
add_slide_title(slide, "Problem Statement")

bullets = [
    "Students and faculty at AdtU lack a dedicated, subject-wise digital workspace for academic collaboration.",
    "Critical academic content (assignments, notes, files) is scattered across WhatsApp, email, and personal drives.",
    "There is no mechanism for real-time, subject-specific communication with proper role enforcement.",
    "Assignment deadlines, grading, and submission tracking are managed manually with no automation.",
    "Announcements posted in generic groups are frequently missed, causing delays and confusion.",
    "Existing LMS platforms are either too complex, unavailable, or not tailored to the AdtU academic structure.",
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(14))

# Highlight box
add_rect(slide, Inches(0.4), Inches(6.1), Inches(12.5), Inches(0.85),
         fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width=Pt(1))
add_textbox(slide, "Core Problem: The absence of a structured, real-time academic collaboration platform "
            "results in information loss, missed deadlines, and poor faculty–student communication.",
            Inches(0.6), Inches(6.18), Inches(12.2), Inches(0.7),
            font_size=Pt(12), bold=True, color=DARK_BLUE, font_name="Calibri")

add_notes(slide, ("Clearly articulate the problem the project solves.\n"
    "Emphasise that this is a real, observed problem at AdtU — not a hypothetical one. "
    "The highlight box at the bottom condenses the entire problem into one concise statement. "
    "This framing justifies why building this system was necessary."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — OBJECTIVES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 1: Introduction")
add_footer_bar(slide)
add_slide_number(slide, 6)
add_slide_title(slide, "Objectives of the Project")

bullets = [
    "Develop a subject-channel-based academic collaboration platform for AdtU students and faculty.",
    "Implement real-time messaging with file attachment support using Socket.IO within each subject channel.",
    "Build a structured assignment management system with deadline tracking, submission, and grading.",
    "Provide a centralised notes and file repository with cloud storage (Cloudinary) per subject.",
    "Implement role-based access control (Student / Class Representative / Faculty / Admin).",
    "Enable faculty to post and auto-distribute announcements with in-app notifications.",
    "Deploy the application on cloud infrastructure (Vercel + Render) for university-wide access.",
    "Design a clean, intuitive UI that requires zero technical training for students or faculty.",
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(13.5))

add_notes(slide, ("Present the specific, measurable objectives of the project.\n"
    "Each objective maps to a concrete implemented feature. These objectives were defined at the "
    "start of the project and guided the development process throughout. Mention that all listed "
    "objectives have been successfully achieved in the final system."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EXISTING SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 2: System Overview")
add_footer_bar(slide)
add_slide_number(slide, 7)
add_slide_title(slide, "Existing System")

left_items = [
    ("Current Tools Used by Students",
     ["WhatsApp groups — primary communication channel",
      "Email — for faculty-to-student announcements",
      "Google Drive / Pen drives — for file sharing",
      "Verbal communication — for assignment details",
      "Google Classroom — used informally, not integrated"]),
]
add_bullet_block(slide, left_items, left=Inches(0.5), top=Inches(1.55),
                 width=Inches(5.8), font_size=Pt(13))

# Vertical divider
add_rect(slide, Inches(6.55), Inches(1.55), Pt(1.5), Inches(5.2), fill_rgb=GOLD)

right_items = [
    ("Limitations of Existing Approach",
     ["No subject-specific separation — all messages mixed in one chat",
      "Files and media get buried under chat history",
      "No assignment tracking, submission, or grading workflow",
      "No role control — anyone can post anything",
      "Real-time notifications are absent or unreliable",
      "No persistent, searchable academic knowledge base"]),
]
add_bullet_block(slide, right_items, left=Inches(6.8), top=Inches(1.55),
                 width=Inches(6.0), font_size=Pt(13))

add_notes(slide, ("Describe the current state of affairs before this project.\n"
    "This slide establishes the baseline. WhatsApp groups, emails, and pen drives are unreliable "
    "and unstructured. Highlight that the fundamental problem is the LACK of structure — not the "
    "absence of communication tools."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PROPOSED SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 2: System Overview")
add_footer_bar(slide)
add_slide_number(slide, 8)
add_slide_title(slide, "Proposed System")

bullets = [
    ("ADTU Collab — Jira for Students",
     ["A full-stack, cloud-hosted academic collaboration platform purpose-built for AdtU.",
      "Organises all academic activity into structured Subject Channels per batch and semester.",
      "Each channel contains: Chat, Notes, Assignments, Files, and Announcements tabs."]),
    ("How It Works",
     ["Students log in using their university roll number and date-of-birth (auto-password).",
      "They are automatically enrolled in all subject channels for their programme & semester.",
      "Real-time Socket.IO messaging allows instant subject-specific conversation.",
      "Faculty create assignments with deadlines; students submit directly in the platform.",
      "All files are uploaded securely to Cloudinary — accessible from anywhere.",
      "Notifications are delivered in-app and in real-time whenever events occur."]),
    ("Deployment",
     ["Frontend hosted on Vercel | Backend hosted on Render | DB on Supabase (PostgreSQL)"]),
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(13))

add_notes(slide, ("Describe the proposed solution clearly and confidently.\n"
    "Emphasise that the system mirrors the subject structure of AdtU — channels are created "
    "per batch, semester, and subject. The DOB-as-password approach removes the need for "
    "password reset flows for students. Highlight cloud deployment for production readiness."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — KEY FEATURES
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 2: System Overview")
add_footer_bar(slide)
add_slide_number(slide, 9)
add_slide_title(slide, "Key Features (Implemented)")

features = [
    ("💬 Real-Time Chat",      "Socket.IO powered messaging with file attachments & typing indicators"),
    ("📚 Notes Repository",    "Faculty/CR-created text notes per subject; searchable and persistent"),
    ("📋 Assignment System",   "Create, submit, and grade assignments with deadline & status tracking"),
    ("📢 Announcements",       "Faculty posts auto-notify all enrolled students in real time"),
    ("🔔 Notifications",       "In-app, real-time notification system for assignments & announcements"),
    ("🔐 Role-Based Access",   "Four roles: Student, Class Representative, Faculty, Admin"),
    ("☁️ Cloud File Storage",  "Cloudinary integration — all files available online, no data loss"),
    ("📅 Calendar View",       "Global deadline calendar showing upcoming assignments across subjects"),
]

col1 = features[:4]
col2 = features[4:]

for i, (icon_title, desc) in enumerate(col1):
    top = Inches(1.55) + i * Inches(1.2)
    add_rect(slide, Inches(0.4), top, Inches(6.0), Inches(1.05),
             fill_rgb=GRAY_LIGHT, line_rgb=LIGHT_BLUE, line_width=Pt(0.5))
    add_textbox(slide, icon_title, Inches(0.55), top + Inches(0.1),
                Inches(5.7), Inches(0.38), font_size=Pt(12), bold=True,
                color=DARK_BLUE, font_name="Calibri")
    add_textbox(slide, desc, Inches(0.55), top + Inches(0.48),
                Inches(5.7), Inches(0.5), font_size=Pt(11),
                color=GRAY_TXT, font_name="Calibri")

for i, (icon_title, desc) in enumerate(col2):
    top = Inches(1.55) + i * Inches(1.2)
    add_rect(slide, Inches(6.9), top, Inches(6.0), Inches(1.05),
             fill_rgb=GRAY_LIGHT, line_rgb=LIGHT_BLUE, line_width=Pt(0.5))
    add_textbox(slide, icon_title, Inches(7.05), top + Inches(0.1),
                Inches(5.7), Inches(0.38), font_size=Pt(12), bold=True,
                color=DARK_BLUE, font_name="Calibri")
    add_textbox(slide, desc, Inches(7.05), top + Inches(0.48),
                Inches(5.7), Inches(0.5), font_size=Pt(11),
                color=GRAY_TXT, font_name="Calibri")

add_notes(slide, ("Walk through each feature card on this slide.\n"
    "Each feature listed here is confirmed as implemented in the codebase. Mention that "
    "the calendar view uses a global assignments endpoint, real-time chat is handled by "
    "Socket.IO events, and cloud storage uses Cloudinary's Node.js SDK with Multer middleware."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 3: System Design")
add_footer_bar(slide)
add_slide_number(slide, 10)
add_slide_title(slide, "System Architecture")

layers = [
    ("Presentation Layer (Frontend)",
     "React 18 + Vite + Tailwind CSS | React Router DOM | Axios | Socket.IO-Client",
     LIGHT_BLUE, DARK_BLUE),
    ("API Layer (Backend)",
     "Node.js + Express.js | REST API (10 route modules) | JWT Authentication | Multer",
     RGBColor(0xD6, 0xF5, 0xE8), RGBColor(0x00, 0x5C, 0x2E)),
    ("Real-Time Layer",
     "Socket.IO Server | Events: message:send, typing:start/stop, notification:send, channel:join/leave",
     RGBColor(0xFF, 0xF3, 0xCD), RGBColor(0x7A, 0x50, 0x00)),
    ("Data Layer (Database)",
     "PostgreSQL on Supabase/Neon | 11 Tables | Indexed for performance | Connection pooling via pg",
     RGBColor(0xF0, 0xD6, 0xF5), RGBColor(0x50, 0x00, 0x7A)),
    ("Storage Layer",
     "Cloudinary CDN | Stores: chat attachments, assignment submissions, notes files",
     RGBColor(0xFF, 0xE5, 0xD6), RGBColor(0x7A, 0x20, 0x00)),
]

for i, (layer_name, details, bg, fg) in enumerate(layers):
    top = Inches(1.58) + i * Inches(1.06)
    add_rect(slide, Inches(0.4), top, Inches(12.5), Inches(0.95),
             fill_rgb=bg, line_rgb=fg, line_width=Pt(0.8))
    add_textbox(slide, layer_name, Inches(0.6), top + Inches(0.05),
                Inches(4.0), Inches(0.38), font_size=Pt(12), bold=True,
                color=fg, font_name="Calibri")
    add_rect(slide, Inches(4.7), top + Inches(0.2), Pt(1), Inches(0.55), fill_rgb=fg)
    add_textbox(slide, details, Inches(4.9), top + Inches(0.12),
                Inches(7.8), Inches(0.7), font_size=Pt(11),
                color=GRAY_TXT, font_name="Calibri")

add_notes(slide, ("Explain the layered architecture of the system.\n"
    "The frontend (React/Vite) communicates with the backend (Express) via REST API and "
    "Socket.IO. The backend handles auth (JWT), file uploads (Multer + Cloudinary), and "
    "DB queries (pg pool). The database is PostgreSQL hosted on Supabase. This five-layer "
    "architecture follows a clean separation of concerns."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — DATABASE DESIGN
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 3: System Design")
add_footer_bar(slide)
add_slide_number(slide, 11)
add_slide_title(slide, "Database Design")

tables = [
    ("faculties",             "id, name, color_code"),
    ("programmes",            "id, faculty_id(FK), name, code, duration_semesters"),
    ("batches",               "id, programme_id(FK), year"),
    ("users",                 "id, name, email, roll_number, dob, password, role, programme_id(FK), batch_year, current_semester"),
    ("channels (subjects)",   "id, batch_id(FK), semester_number, subject_name, teacher_id(FK)"),
    ("enrollments",           "id, user_id(FK), channel_id(FK) — UNIQUE pair"),
    ("messages",              "id, channel_id(FK), sender_id(FK), content, file_url, is_pinned"),
    ("files",                 "id, channel_id(FK), uploaded_by(FK), file_name, file_url, file_type, file_size"),
    ("notes",                 "id, channel_id(FK), created_by(FK), title, content"),
    ("assignments",           "id, channel_id(FK), created_by(FK), title, due_date, max_marks"),
    ("assignment_submissions","id, assignment_id(FK), student_id(FK), file_url, status, marks, feedback"),
    ("announcements",         "id, channel_id(FK), user_id(FK), title, content, is_important"),
    ("notifications",         "id, user_id(FK), type, title, message, is_read, ref_id"),
]

# Two columns
col_break = 7
for i, (tbl, cols) in enumerate(tables):
    if i < col_break:
        lft, top = Inches(0.4), Inches(1.55) + i * Inches(0.73)
        w = Inches(6.1)
    else:
        lft, top = Inches(6.85), Inches(1.55) + (i - col_break) * Inches(0.73)
        w = Inches(6.1)
    bg = GRAY_LIGHT if i % 2 == 0 else WHITE
    add_rect(slide, lft, top, w, Inches(0.65),
             fill_rgb=bg, line_rgb=LIGHT_BLUE, line_width=Pt(0.5))
    add_textbox(slide, tbl, lft + Inches(0.1), top + Inches(0.04),
                Inches(2.1), Inches(0.35), font_size=Pt(10), bold=True,
                color=DARK_BLUE, font_name="Calibri")
    add_textbox(slide, cols, lft + Inches(2.3), top + Inches(0.04),
                w - Inches(2.4), Inches(0.55), font_size=Pt(9),
                color=GRAY_TXT, font_name="Calibri")

add_notes(slide, ("Explain the PostgreSQL database schema.\n"
    "There are 13 tables. The hierarchical structure is: Faculty → Programme → Batch → Channel. "
    "Users are linked to a programme and batch. Enrollments connect users to channels. "
    "All activity tables (messages, files, notes, assignments, notifications) reference channels "
    "and users via foreign keys. Performance indexes are applied on all FK columns."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — WORKFLOW / PROCESS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 3: System Design")
add_footer_bar(slide)
add_slide_number(slide, 12)
add_slide_title(slide, "System Workflow / Process Flow")

steps = [
    ("1", "Student Registration",
     "Student signs up using roll number, DOB (password), and selects programme & batch."),
    ("2", "Auto Enrolment",
     "Backend auto-enrols the student in all subject channels for their programme and current semester."),
    ("3", "Dashboard Access",
     "Student lands on a personalised dashboard showing all enrolled subject channels as cards."),
    ("4", "Subject Channel",
     "Student opens a subject to access Chat, Notes, Assignments, and Files tabs."),
    ("5", "Real-Time Chat",
     "Messages sent via Socket.IO; persisted to PostgreSQL; delivered instantly to all channel members."),
    ("6", "Assignment Lifecycle",
     "Faculty creates → Students submit (file upload to Cloudinary) → Faculty grades → Notification sent."),
    ("7", "Notifications",
     "In-app notifications triggered for: new assignment, graded submission, new announcement."),
]

step_h = Inches(0.73)
for i, (num, title, desc) in enumerate(steps):
    top = Inches(1.55) + i * step_h
    add_rect(slide, Inches(0.4), top, Inches(0.6), step_h - Pt(3),
             fill_rgb=DARK_BLUE)
    add_textbox(slide, num, Inches(0.4), top + Inches(0.17),
                Inches(0.6), Inches(0.35), font_size=Pt(14), bold=True,
                color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(slide, title, Inches(1.1), top + Inches(0.04),
                Inches(2.8), Inches(0.35), font_size=Pt(12), bold=True,
                color=DARK_BLUE, font_name="Calibri")
    add_textbox(slide, desc, Inches(4.0), top + Inches(0.04),
                Inches(9.0), Inches(0.6), font_size=Pt(11.5),
                color=GRAY_TXT, font_name="Calibri")
    if i < len(steps) - 1:
        add_rect(slide, Inches(0.65), top + step_h - Pt(2), Pt(2),
                 Pt(4), fill_rgb=GOLD)

add_notes(slide, ("Walk through the end-to-end user journey step by step.\n"
    "Start from registration, explain auto-enrolment (the system reads programme and batch "
    "to assign channels automatically), then show how a student uses the system daily. "
    "Emphasise the assignment lifecycle as a key workflow that replaces manual processes."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — TECHNOLOGIES USED
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 4: Implementation")
add_footer_bar(slide)
add_slide_number(slide, 13)
add_slide_title(slide, "Technologies Used")

tech_groups = [
    ("Frontend",  ["React 18", "Vite 5", "Tailwind CSS 3", "React Router DOM v6",
                   "Axios 1.7", "Socket.IO-Client 4.7", "date-fns 3.6"]),
    ("Backend",   ["Node.js", "Express.js 4.18", "Socket.IO 4.7",
                   "Multer 1.4 (file handling)", "Multer-Storage-Cloudinary",
                   "bcryptjs 2.4 (password hashing)", "jsonwebtoken 9.0",
                   "express-async-errors"]),
    ("Database",  ["PostgreSQL (via pg 8.11)", "Supabase / Neon (cloud hosted)",
                   "Connection pooling", "13 relational tables"]),
    ("Cloud/DevOps", ["Cloudinary (file & media CDN)", "Vercel (frontend deployment)",
                      "Render (backend deployment)", "dotenv (env management)",
                      "Nodemon (development)", "Git + GitHub (version control)"]),
]

col_w = Inches(6.0)
for gi, (group, techs) in enumerate(tech_groups):
    col = gi % 2
    row = gi // 2
    lft = Inches(0.4) + col * Inches(6.7)
    top = Inches(1.6) + row * Inches(2.65)
    add_rect(slide, lft, top, col_w, Inches(0.38), fill_rgb=DARK_BLUE)
    add_textbox(slide, group.upper(), lft + Inches(0.1), top + Inches(0.04),
                col_w - Inches(0.2), Inches(0.3),
                font_size=Pt(11), bold=True, color=WHITE, font_name="Calibri")
    for ti, t in enumerate(techs):
        add_textbox(slide, f"•  {t}",
                    lft + Inches(0.15), top + Inches(0.42) + ti * Inches(0.29),
                    col_w - Inches(0.25), Inches(0.27),
                    font_size=Pt(11), color=GRAY_TXT, font_name="Calibri")

add_notes(slide, ("List and briefly explain each technology.\n"
    "Frontend uses React + Vite for fast development, Tailwind for utility-first styling, "
    "and React Router for SPA navigation. Backend uses Express as the REST framework, "
    "Socket.IO for WebSocket real-time features, and JWT for stateless authentication. "
    "All files go through Multer to Cloudinary. Database is PostgreSQL with Supabase for hosting."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — UI SCREENS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 4: Implementation")
add_footer_bar(slide)
add_slide_number(slide, 14)
add_slide_title(slide, "User Interface — Key Screens")

screens = [
    ("Login / Signup Page",
     "Roll number + DOB authentication. Student selects programme & batch on signup."),
    ("Dashboard",
     "Grid of subject channel cards showing programme, semester, and teacher info."),
    ("Subject Page — Dashboard Tab",
     "Stats row (files, pending tasks, subjects) + pinned announcement + assignments list + live chat panel."),
    ("Subject Page — Chat Tab",
     "Real-time Socket.IO chat with message bubbles, file attachments, role badges (Teacher / CR)."),
    ("Subject Page — Assignments Tab",
     "Assignment cards with due-date status pills (Overdue / Due Soon / Completed) + Submit Work button."),
    ("Subject Page — Files Tab",
     "File cards with type icons, description, size, and direct download link from Cloudinary."),
    ("Notifications Page",
     "In-app notification feed grouped by type (assignment, announcement) with read/unread states."),
    ("Calendar Page",
     "Month view of all assignment deadlines across enrolled subjects with event indicators."),
]

for i, (screen, desc) in enumerate(screens):
    col = i % 2
    row = i // 2
    lft = Inches(0.4) + col * Inches(6.5)
    top = Inches(1.58) + row * Inches(1.32)
    add_rect(slide, lft, top, Inches(6.2), Inches(1.2),
             fill_rgb=GRAY_LIGHT, line_rgb=MID_BLUE, line_width=Pt(0.5))
    # Placeholder icon
    add_rect(slide, lft + Inches(0.1), top + Inches(0.08),
             Inches(0.7), Inches(1.0), fill_rgb=LIGHT_BLUE)
    add_textbox(slide, "📸", lft + Inches(0.1), top + Inches(0.3),
                Inches(0.7), Inches(0.4),
                font_size=Pt(16), align=PP_ALIGN.CENTER, font_name="Calibri")
    add_textbox(slide, screen, lft + Inches(0.9), top + Inches(0.08),
                Inches(5.2), Inches(0.35), font_size=Pt(12), bold=True,
                color=DARK_BLUE, font_name="Calibri")
    add_textbox(slide, desc, lft + Inches(0.9), top + Inches(0.44),
                Inches(5.2), Inches(0.7), font_size=Pt(10.5),
                color=GRAY_TXT, font_name="Calibri")

add_notes(slide, ("Describe each screen of the application.\n"
    "Replace the 📸 placeholder icons with actual application screenshots during final preparation. "
    "Walk through each screen logically: start with login, move to dashboard, then into a subject "
    "channel showing the different tabs. Highlight the split-panel dashboard design (stats + chat)."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — METHODOLOGY
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 4: Implementation")
add_footer_bar(slide)
add_slide_number(slide, 15)
add_slide_title(slide, "Development Methodology")

bullets = [
    ("Approach: Agile (Iterative & Incremental)",
     ["Development was divided into iterative phases with continuous review and refinement.",
      "Features were built, tested, and integrated module by module rather than all at once."]),
    ("Phase 1 — Foundation",
     ["Database schema design (PostgreSQL / Supabase)",
      "Authentication system: JWT-based login, role assignment, and roll-number verification",
      "Core Express.js server setup with CORS and middleware pipeline"]),
    ("Phase 2 — Core Features",
     ["Subject channel system with auto-enrolment logic",
      "Real-time chat via Socket.IO (messages, typing indicators, file attachments)",
      "Assignment creation, submission (Cloudinary upload), and grading workflow"]),
    ("Phase 3 — Enrichment & Deployment",
     ["Notes repository, file manager, announcement system, in-app notifications",
      "Calendar view and global All Files page",
      "Production deployment: Vercel (frontend) + Render (backend) + Supabase (DB)"]),
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(13))

add_notes(slide, ("Explain the development approach and phases.\n"
    "The project followed an Agile-inspired iterative model. Phase 1 was all about the "
    "data foundation and security. Phase 2 introduced the core real-time features. "
    "Phase 3 polished the system and made it production-ready. This phased approach "
    "ensured each layer was stable before building on top of it."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — RESULTS / OUTPUT
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 5: Results & Discussion")
add_footer_bar(slide)
add_slide_number(slide, 16)
add_slide_title(slide, "Results & System Output")

bullets = [
    ("Successfully Implemented Features",
     ["Real-time subject-channel chat with file attachments fully functional via Socket.IO",
      "Complete assignment lifecycle: create → submit → grade → notification",
      "Role-based access: Students, Class Representatives, Faculty, Admin — all distinct",
      "Cloudinary file storage operational for messages, assignments, and notes",
      "In-app notification system with real-time delivery and read/unread states",
      "Calendar page aggregating all assignment deadlines across enrolled subjects",
      "Production deployment live on Vercel + Render with Supabase PostgreSQL"]),
    ("Backend API Coverage",
     ["10 route modules: auth, channels, messages, files, assignments, submissions,",
      "announcements, notifications, notes, enrollments",
      "Health endpoint verified: /api/health returns 200 with status message"]),
    ("Data Integrity",
     ["13 PostgreSQL tables with FK constraints, CASCADE deletes, and UNIQUE constraints",
      "Performance indexes on all high-query columns (messages, enrollments, files)"]),
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(13))

add_notes(slide, ("Present the working results of the system.\n"
    "This slide demonstrates that the project is not just conceptual — it is a fully working, "
    "deployed application. Mention specific figures: 10 API route modules, 13 database tables, "
    "4 user roles, 8 UI screens. If possible, do a live demo after this slide."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — CHALLENGES FACED
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 5: Results & Discussion")
add_footer_bar(slide)
add_slide_number(slide, 17)
add_slide_title(slide, "Challenges Faced")

bullets = [
    ("Real-Time Socket.IO Integration",
     ["Synchronising Socket.IO events with REST API state without duplicate messages",
      "Handling socket room management (channel:join / channel:leave) across React re-renders",
      "Solution: Used useRef to persist socket instance across renders; cleaned up on unmount"]),
    ("Database Schema Design",
     ["Modelling the Faculty → Programme → Batch → Channel hierarchy correctly",
      "Ensuring auto-enrolment logic correctly maps students to channels without duplication",
      "Solution: UNIQUE(user_id, channel_id) constraint in enrollments table"]),
    ("File Upload with Cloudinary",
     ["Configuring Multer + multer-storage-cloudinary pipeline in a streaming manner",
      "Handling large files and CORS issues between Render (backend) and Vercel (frontend)",
      "Solution: Configured explicit CORS with regex for Vercel preview URLs"]),
    ("JWT Authentication Across Roles",
     ["Ensuring role-based middleware correctly restricts endpoints per user role",
      "Managing token refresh and expiry without breaking user session",
      "Solution: Decoded token in auth middleware and attached user object to req"]),
    ("Deployment Configuration",
     ["Supabase transaction pooler credentials differing from session pooler",
      "Environment variable management across local dev, Render, and Vercel",
      "Solution: Used .env.example template and platform-specific env variable dashboards"]),
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(12))

add_notes(slide, ("Discuss the real technical challenges encountered and how they were resolved.\n"
    "These are genuine issues faced during development — not hypothetical. The CORS/Vercel issue, "
    "the Supabase pooler configuration, and socket cleanup in React are all real documented problems "
    "that were systematically resolved. This shows technical depth and problem-solving ability."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — CONCLUSION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 6: Conclusion & Future Scope")
add_footer_bar(slide)
add_slide_number(slide, 18)
add_slide_title(slide, "Conclusion")

bullets = [
    "The project successfully delivers a fully functional, cloud-deployed academic collaboration platform tailored for AdtU students and faculty.",
    "It addresses the core problem of scattered, unstructured academic communication by providing dedicated subject channels with structured workflows.",
    "Key achievements: real-time messaging, role-based access control, cloud file storage, assignment lifecycle management, and in-app notifications — all production-ready.",
    "The system is built on a modern, industry-standard technology stack (React, Node.js, PostgreSQL, Socket.IO, Cloudinary) ensuring maintainability and scalability.",
    "The platform has the potential to replace informal tools (WhatsApp groups, email chains) as the primary academic communication medium at Assam down town University.",
    "The project demonstrates the practical application of full-stack web development, cloud deployment, real-time systems, and database design concepts learned during BSc-IT (MAIS).",
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(13.5))

add_notes(slide, ("Wrap up the project with a strong, confident conclusion.\n"
    "Emphasise that every objective stated at the beginning has been achieved. The system is not "
    "a prototype — it is deployed and functional. Highlight that this project bridges academic "
    "learning with real-world software engineering practice."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — FUTURE SCOPE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide, "Chapter 6: Conclusion & Future Scope")
add_footer_bar(slide)
add_slide_number(slide, 19)
add_slide_title(slide, "Future Scope")

bullets = [
    ("Short-Term Enhancements",
     ["Mobile application (React Native) for iOS and Android with push notifications",
      "Rich-text message editor (markdown support) in the chat interface",
      "Advanced assignment grading with rubric-based scoring and batch feedback"]),
    ("Medium-Term Expansions",
     ["AI-powered smart search across all notes, files, and messages using embeddings",
      "Automated attendance tracking integrated with subject channels",
      "Faculty analytics dashboard — submission rates, active students, engagement metrics"]),
    ("Long-Term Vision",
     ["University-wide adoption — multi-campus, multi-department scaling",
      "Integration with AdtU's official ERP/SIS (Student Information System)",
      "Video lecture integration (Zoom/Google Meet embed) within subject channels",
      "Peer-to-peer collaborative note editing (Google Docs-style real-time editing)"]),
]
add_bullet_block(slide, bullets, top=Inches(1.55), font_size=Pt(13))

add_notes(slide, ("Discuss planned and potential future improvements.\n"
    "Frame future scope as a roadmap — short-term (3–6 months), medium-term (6–12 months), "
    "and long-term (1–2 years). The mobile app and AI search are highest priority. "
    "ERP integration would be the most transformative for the university."))


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — REFERENCES + THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_header_bar(slide)
add_footer_bar(slide)
add_slide_number(slide, 20)
add_slide_title(slide, "References")

refs = [
    "[1]  React Documentation — https://react.dev",
    "[2]  Node.js Documentation — https://nodejs.org/en/docs",
    "[3]  Express.js Guide — https://expressjs.com/en/guide",
    "[4]  Socket.IO Documentation — https://socket.io/docs/v4",
    "[5]  PostgreSQL Documentation — https://www.postgresql.org/docs",
    "[6]  Supabase Documentation — https://supabase.com/docs",
    "[7]  Cloudinary Node.js SDK — https://cloudinary.com/documentation/node_integration",
    "[8]  Vite.js Documentation — https://vitejs.dev/guide",
    "[9]  JSON Web Tokens (JWT) — https://jwt.io/introduction",
    "[10] Multer Middleware — https://github.com/expressjs/multer",
    "[11] Tailwind CSS Documentation — https://tailwindcss.com/docs",
    "[12] Pressman, R.S. — Software Engineering: A Practitioner's Approach (8th Ed.), McGraw-Hill",
]

for i, ref in enumerate(refs):
    add_textbox(slide, ref,
                Inches(0.5), Inches(1.55) + i * Inches(0.37),
                Inches(12.3), Inches(0.35),
                font_size=Pt(11), color=GRAY_TXT, font_name="Calibri")

# Thank You banner
add_rect(slide, Inches(2.0), Inches(6.05), Inches(9.3), Inches(0.9),
         fill_rgb=DARK_BLUE, line_rgb=GOLD, line_width=Pt(1.5))
add_textbox(slide, "Thank You  |  Questions & Discussion",
            Inches(2.0), Inches(6.13), Inches(9.3), Inches(0.7),
            font_size=Pt(18), bold=True, color=WHITE,
            align=PP_ALIGN.CENTER, font_name="Calibri")

add_notes(slide, ("Closing slide — thank the audience and invite questions.\n"
    "Reference [12] (Pressman) represents the software engineering methodology foundation. "
    "All other references are official documentation of the tools actually used in the project. "
    "Invite evaluators to ask questions about any aspect of the implementation."))


# ── SAVE ───────────────────────────────────────────────────────────────────────
output_path = r"c:\Users\Sonux\Desktop\major project\Jira_For_Students_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Total slides: {len(prs.slides)}")
